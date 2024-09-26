

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

# Create a new PowerPoint presentation
prs = Presentation()

# Add a blank slide layout
slide_layout = prs.slide_layouts[5]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Define starting position and sizes
left = Inches(1)
top = Inches(1)
width = Inches(2)
height = Inches(1)

# Add the main process (HR Recruiter Website)
hr_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
hr_shape.text = "HR Recruiter Website (Main System)"

# Position next step (Technical Support)
top += Inches(1.5)
support_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
support_shape.text = "Technical Support\n- User Assistance\n- Support Tickets\n- Live Chat Support"

# Add connector between main system and technical support
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, hr_shape.left + width/2, hr_shape.top + height,
                                       support_shape.left + width/2, support_shape.top)
connector.line.width = Inches(0.05)

# Position next step (Bug Fixing and Issue Resolution)
top += Inches(1.5)
bug_fix_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
bug_fix_shape.text = "Bug Fixing & Issue Resolution\n- Monitoring & Debugging\n- Bug Fix & Patch Deployment\n- Regression Testing"

# Add connector between technical support and bug fixing
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, support_shape.left + width/2, support_shape.top + height,
                                       bug_fix_shape.left + width/2, bug_fix_shape.top)
connector.line.width = Inches(0.05)

# Position next step (Server & Website Maintenance)
top += Inches(1.5)
server_maintenance_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
server_maintenance_shape.text = "Server & Website Maintenance\n- Uptime Monitoring\n- Server Health Checks\n- Load Balancing"

# Add connector between bug fixing and server maintenance
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, bug_fix_shape.left + width/2, bug_fix_shape.top + height,
                                       server_maintenance_shape.left + width/2, server_maintenance_shape.top)
connector.line.width = Inches(0.05)

# Position next step (Data Backup & Security)
top += Inches(1.5)
data_backup_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
data_backup_shape.text = "Data Backup & Security\n- Database Backups\n- Security Monitoring\n- User Authentication"

# Add connector between server maintenance and data backup
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, server_maintenance_shape.left + width/2, server_maintenance_shape.top + height,
                                       data_backup_shape.left + width/2, data_backup_shape.top)
connector.line.width = Inches(0.05)

# Position next step (User Account Management)
top += Inches(1.5)
user_account_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
user_account_shape.text = "User Account Management\n- New User Onboarding\n- Account Support\n- Inactivity Monitoring"

# Add connector between data backup and user account management
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, data_backup_shape.left + width/2, data_backup_shape.top + height,
                                       user_account_shape.left + width/2, user_account_shape.top)
connector.line.width = Inches(0.05)

# Position next step (Feature Updates & Improvements)
top += Inches(1.5)
feature_updates_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
feature_updates_shape.text = "Feature Updates & Improvements\n- New Feature Rollouts\n- UX Testing & Feedback"

# Add connector between user account management and feature updates
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, user_account_shape.left + width/2, user_account_shape.top + height,
                                       feature_updates_shape.left + width/2, feature_updates_shape.top)
connector.line.width = Inches(0.05)

# Position next step (Performance Optimization)
top += Inches(1.5)
performance_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
performance_shape.text = "Performance Optimization\n- System Performance Monitoring\n- Speed Optimization\n- Scalability Planning"

# Add connector between feature updates and performance optimization
connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, feature_updates_shape.left + width/2, feature_updates_shape.top + height,
                                       performance_shape.left + width/2, performance_shape.top)
connector.line.width = Inches(0.05)

# Save the PowerPoint presentation
prs.save('HR_Recruiter_Flowchart.pptx')
